"""
AI Study Assistant - FastAPI Backend
Features: RAG, PDF upload, Quiz generation, Summarization, Session tracking
"""

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
import os, io, json, uuid, hashlib
from dotenv import load_dotenv

load_dotenv()

app = FastAPI(title="AI Study Assistant API")

# ── Provider setup ────────────────────────────────────────────────────────────
PROVIDER = os.getenv("AI_PROVIDER", "").lower()

if PROVIDER == "openai":
    from openai import OpenAI
    ai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
elif PROVIDER == "gemini":
    import google.generativeai as genai
    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
    ai_client = None  # handled per-call
elif PROVIDER == "anthropic":
    import anthropic
    ai_client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
else:
    raise RuntimeError("AI_PROVIDER not set. Add AI_PROVIDER=openai|gemini|anthropic to .env")

# ── RAG: In-memory vector store ───────────────────────────────────────────────
import re

document_store: dict[str, dict] = {}  # session_id -> {chunks, filename}


def extract_text_from_pdf(file_bytes: bytes) -> str:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
    return "\n".join(text)


def chunk_text(text: str, size: int = 500, overlap: int = 50) -> list[str]:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunk = " ".join(words[i:i + size])
        if chunk.strip():
            chunks.append(chunk)
        i += size - overlap
    return chunks


def simple_embed(text: str) -> list[float]:
    """Lightweight TF-style embedding using word frequency (no external model needed)."""
    words = re.findall(r'\w+', text.lower())
    freq: dict[str, int] = {}
    for w in words:
        freq[w] = freq.get(w, 0) + 1
    total = sum(freq.values()) or 1
    return {w: c / total for w, c in freq.items()}


def cosine_sim(a: dict, b: dict) -> float:
    keys = set(a) & set(b)
    if not keys:
        return 0.0
    dot = sum(a[k] * b[k] for k in keys)
    mag_a = sum(v ** 2 for v in a.values()) ** 0.5
    mag_b = sum(v ** 2 for v in b.values()) ** 0.5
    return dot / (mag_a * mag_b) if mag_a and mag_b else 0.0


def retrieve_context(session_id: str, query: str, top_k: int = 4) -> str:
    if session_id not in document_store:
        return ""
    q_vec = simple_embed(query)
    chunks_data = document_store[session_id]["chunks"]
    scored = [(cosine_sim(q_vec, c["vec"]), c["text"]) for c in chunks_data]
    scored.sort(reverse=True)
    return "\n\n".join(text for _, text in scored[:top_k])


# ── AI call helpers ───────────────────────────────────────────────────────────
SUBJECTS = ["Mathematics", "Physics", "Chemistry", "Biology",
            "Computer Science", "History", "Economics", "Literature"]


def call_ai_full(system: str, user: str) -> str:
    if PROVIDER == "openai":
        r = ai_client.chat.completions.create(
            model="gpt-4o-mini", max_tokens=1500,
            messages=[{"role": "system", "content": system}, {"role": "user", "content": user}]
        )
        return r.choices[0].message.content

    elif PROVIDER == "gemini":
        m = genai.GenerativeModel(model_name="gemini-2.0-flash", system_instruction=system)
        return m.generate_content(user).text

    elif PROVIDER == "anthropic":
        r = ai_client.messages.create(
            model="claude-haiku-4-5", max_tokens=1500,
            system=system, messages=[{"role": "user", "content": user}]
        )
        return r.content[0].text


def stream_ai(system: str, user: str):
    if PROVIDER == "openai":
        with ai_client.chat.completions.create(
            model="gpt-4o-mini", max_tokens=800, stream=True,
            messages=[{"role": "system", "content": system}, {"role": "user", "content": user}]
        ) as s:
            for chunk in s:
                t = chunk.choices[0].delta.content
                if t:
                    yield t

    elif PROVIDER == "gemini":
        m = genai.GenerativeModel(model_name="gemini-2.0-flash", system_instruction=system)
        for chunk in m.generate_content(user, stream=True):
            if chunk.text:
                yield chunk.text

    elif PROVIDER == "anthropic":
        with ai_client.messages.stream(
            model="claude-haiku-4-5", max_tokens=800,
            system=system, messages=[{"role": "user", "content": user}]
        ) as s:
            for text in s.text_stream:
                yield text


# ── Models ────────────────────────────────────────────────────────────────────
class QuestionRequest(BaseModel):
    subject: str
    question: str
    level: str = "Undergraduate"
    session_id: str = ""


class QuizRequest(BaseModel):
    session_id: str
    subject: str
    level: str = "Undergraduate"
    num_questions: int = 5


class QuizAnswer(BaseModel):
    session_id: str
    question: str
    selected: str
    correct: str
    subject: str


# ── Routes ────────────────────────────────────────────────────────────────────
@app.get("/")
def root():
    return {"status": "running", "provider": PROVIDER}


@app.get("/subjects")
def get_subjects():
    return {"subjects": SUBJECTS}


@app.post("/upload")
async def upload_pdf(file: UploadFile = File(...)):
    filename = file.filename.lower()
    if not (filename.endswith(".pdf") or filename.endswith(".pptx")):
        raise HTTPException(400, "Only PDF and PPTX files are supported.")

    content = await file.read()

    if filename.endswith(".pdf"):
        text = extract_text_from_pdf(content)
    else:
        text = extract_text_from_pptx(content)

    if not text.strip():
        raise HTTPException(400, "Could not extract text from file.")

    chunks = chunk_text(text)
    session_id = hashlib.md5(content).hexdigest()[:12]
    document_store[session_id] = {
        "filename": file.filename,
        "chunks": [{"text": c, "vec": simple_embed(c)} for c in chunks],
        "full_text": text[:8000]
    }
    return {"session_id": session_id, "filename": file.filename, "chunks": len(chunks)}


@app.post("/summarize")
def summarize(session_id: str, subject: str = "", level: str = "Undergraduate"):
    if session_id not in document_store:
        raise HTTPException(404, "Session not found. Please upload a PDF first.")
    doc = document_store[session_id]
    system = f"""You are an academic study assistant. Summarize lecture slides clearly for a {level} student.
Structure your summary as:
1. 📌 Main Topics (bullet points)
2. 🔑 Key Concepts (brief definitions)
3. 💡 Important Takeaways"""
    user = f"Summarize this lecture content:\n\n{doc['full_text']}"
    summary = call_ai_full(system, user)
    return {"summary": summary, "filename": doc["filename"]}


@app.post("/quiz")
def generate_quiz(req: QuizRequest):
    if req.session_id and req.session_id in document_store:
        context = document_store[req.session_id]["full_text"][:4000]
        source = "lecture slides"
    else:
        context = ""
        source = f"{req.subject} at {req.level} level"

    system = """You are a quiz generator. Return ONLY a valid JSON array, no markdown, no explanation.
Each item must have: question (string), options (array of 4 strings), answer (one of the 4 strings exactly)."""

    user = f"""Generate {req.num_questions} multiple choice questions about {source}.
Level: {req.level}
{"Context:\n" + context if context else ""}

Return only JSON like: [{{"question":"...","options":["A","B","C","D"],"answer":"A"}}]"""

    raw = call_ai_full(system, user)
    raw = raw.strip().removeprefix("```json").removeprefix("```").removesuffix("```").strip()
    try:
        quiz = json.loads(raw)
        return {"quiz": quiz}
    except Exception:
        raise HTTPException(500, "Failed to parse quiz. Try again.")


@app.post("/ask")
def ask_question(req: QuestionRequest):
    if not req.question.strip():
        raise HTTPException(400, "Question cannot be empty.")

    context = retrieve_context(req.session_id, req.question) if req.session_id else ""

    system = f"""You are a strict study assistant for {req.subject} at {req.level} level.
ONLY answer questions related to {req.subject}. Refuse anything off-topic.
If lecture context is provided, prioritize it in your answer.
Structure: direct answer → explanation → example."""

    user = (f"Lecture Context:\n{context}\n\n" if context else "") + \
           f"Subject: {req.subject} | Level: {req.level}\nQuestion: {req.question}"

    return StreamingResponse(stream_ai(system, user), media_type="text/plain")
